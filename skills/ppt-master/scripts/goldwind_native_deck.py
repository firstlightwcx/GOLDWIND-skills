#!/usr/bin/env python3
"""Build a Goldwind-standard PPTX using native PowerPoint layouts.

This path intentionally avoids SVG-to-PPTX for Goldwind decks. It starts from a
real Goldwind PPTX base file so master/layout elements such as logo, side rail,
rotated copyright text, page-number placeholders, and background artwork are
kept as PowerPoint-native objects. Generated content remains editable.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from copy import deepcopy
from datetime import date
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:  # pragma: no cover - environment guard
    raise SystemExit("python-pptx is required. Install with: python3 -m pip install python-pptx") from exc


PRIMARY = RGBColor(0x06, 0x33, 0x60)
PRIMARY_DEEP = RGBColor(0x02, 0x31, 0x62)
BODY = RGBColor(0x51, 0x51, 0x51)
MUTED = RGBColor(0x84, 0x84, 0x84)
TEAL = RGBColor(0x03, 0x4F, 0x75)
GRAY = RGBColor(0xE9, 0xEA, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDD, 0xDD, 0xDD)

SLIDE_W = 13.333333
SLIDE_H = 7.5
TEXT_FONT = "微软雅黑"
LATIN_FONT = "Arial"
COPYRIGHT_TEXT = "GOLDWIND SCIENCE"


def skill_dir() -> Path:
    return Path(__file__).resolve().parent.parent


def template_dir() -> Path:
    return skill_dir() / "templates" / "layouts" / "金风通用模板"


def default_base() -> Path:
    return template_dir() / "goldwind_native_base.pptx"


def default_toc_image() -> Path:
    return template_dir() / "toc_wind_left.png"


def inches(value: float):
    return Inches(value)


def clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def get_layout(prs: Presentation, name_fragment: str):
    for layout in prs.slide_layouts:
        if name_fragment in layout.name:
            return layout
    raise SystemExit(f"Cannot find slide layout containing: {name_fragment}")


def find_copyright_element(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and COPYRIGHT_TEXT in shape.text.upper() and "©" in shape.text:
                return deepcopy(shape._element)  # noqa: SLF001 - cloning original PPT XML preserves native rendering.
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if getattr(shape, "has_text_frame", False) and COPYRIGHT_TEXT in shape.text.upper() and "©" in shape.text:
                return deepcopy(shape._element)  # noqa: SLF001
    return None


def clone_shape_to_slide(slide, element) -> None:
    if element is None:
        add_copyright(slide)
        return
    cloned = deepcopy(element)
    max_id = max((shape.shape_id for shape in slide.shapes), default=1)
    for node in cloned.iter():
        if node.tag.endswith("}cNvPr"):
            node.set("id", str(max_id + 1))
            node.set("name", "Goldwind Side Copyright")
            break
    slide.shapes._spTree.insert_element_before(cloned, "p:extLst")  # noqa: SLF001


def set_text_style(run, size: float, color=BODY, bold: bool = False, font: str = TEXT_FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text_box(slide, x, y, w, h, text="", size=18, color=BODY, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text_style(run, size, color, bold)
    return shape


def top_title_placeholder(slide):
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False) or not getattr(shape, "has_text_frame", False):
            continue
        if shape.top <= Inches(1.0) and shape.left <= Inches(1.0):
            return shape
    return None


def remove_top_title_placeholder(slide) -> None:
    shape = top_title_placeholder(slide)
    if shape is not None:
        shape._element.getparent().remove(shape._element)  # noqa: SLF001 - remove unused placeholder prompt.


def set_title_placeholder(slide, text: str) -> None:
    shape = top_title_placeholder(slide)
    if shape is None:
        text_box(slide, 0.763, 0.276, 6.8, 0.50, text, 20, PRIMARY, True)
        return
    shape.left = inches(0.763)
    shape.top = inches(0.276)
    shape.width = inches(6.8)
    shape.height = inches(0.50)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_text_style(run, 20, PRIMARY, True)


def rect_text(slide, x, y, w, h, text="", fill=WHITE, line=BORDER, size=16, color=BODY, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = inches(0.16)
    tf.margin_right = inches(0.16)
    tf.margin_top = inches(0.10)
    tf.margin_bottom = inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_text_style(run, size, color, bold)
    return shape


def add_copyright(slide) -> None:
    # Match the native Goldwind report samples: rotated rectangle, not SVG matrix text.
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(-2.376), inches(3.371), inches(5.512), inches(0.76))
    shape.rotation = 270
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "© GOLDWIND SCIENCE & TECHNOLOGY CO., LTD."
    set_text_style(run, 8, BODY, False, LATIN_FONT)


def add_accent(slide, y=2.67):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(0.76), inches(y), inches(0.028), inches(0.787))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_cover(slide, spec: dict) -> None:
    # Side copyright, logo, and wave/background are inherited from the native layout.
    add_accent(slide, 2.67)
    title = spec.get("title") or "金风科技专题汇报"
    subtitle = spec.get("subtitle") or ""
    author = spec.get("author") or "汇报人 - Codex"
    report_date = spec.get("date") or date.today().strftime("%Y-%m-%d")
    text_box(slide, 0.60, 1.55, 12.2, 0.62, title, 33, PRIMARY, True, PP_ALIGN.CENTER)
    if subtitle:
        text_box(slide, 0.60, 2.18, 12.2, 0.48, subtitle, 24, PRIMARY, True, PP_ALIGN.CENTER)
    text_box(slide, 2.15, 4.24, 9.44, 0.39, author, 18, BODY, False, PP_ALIGN.CENTER)
    text_box(slide, 6.17, 4.71, 2.35, 0.38, report_date, 18, BODY, False, PP_ALIGN.CENTER)


def add_toc(slide, items: list[str], toc_image: Path) -> None:
    remove_top_title_placeholder(slide)
    slide.shapes.add_picture(str(toc_image), inches(0), inches(-0.006), width=inches(6.92), height=inches(7.506))
    box = slide.shapes.add_textbox(inches(7.653), inches(0.187), inches(4.899), inches(6.614))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "目录"
    set_text_style(r, 28, BODY, True)
    for idx, title in enumerate(items[:4], start=1):
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(20)
        p.line_spacing = 2.0
        p.text = f"{idx}.  {title}"
        for run in p.runs:
            set_text_style(run, 20, BODY, True)


def add_content(slide, page: dict, section_num: str) -> None:
    title = page.get("title") or "页面标题"
    lead = page.get("lead") or ""
    bullets = page.get("bullets") or []
    source = page.get("source") or ""
    set_title_placeholder(slide, f"{section_num}「{title}」")
    if lead:
        rect_text(slide, 0.763, 1.04, 12.107, 0.86, lead, fill=GRAY, line=GRAY, size=16, color=BODY, bold=True)
    y = 2.18
    for i, item in enumerate(bullets[:4]):
        if isinstance(item, dict):
            heading = item.get("heading", "")
            body = item.get("body", "")
        else:
            heading = str(item)
            body = ""
        x = 0.92 + (i % 2) * 6.02
        yy = y + (i // 2) * 1.72
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x), inches(yy), inches(5.55), inches(1.28))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x), inches(yy), inches(5.55), inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()
        text_box(slide, x + 0.18, yy + 0.26, 5.15, 0.25, heading, 15, PRIMARY, True)
        if body:
            text_box(slide, x + 0.18, yy + 0.62, 5.15, 0.45, body, 12.5, BODY, False)
    if source:
        text_box(slide, 0.84, 7.10, 8.0, 0.18, f"Source: {source}", 8.5, RGBColor(0xA6, 0xA9, 0xAC), False)


def add_ending(slide, spec: dict) -> None:
    # Side copyright, logo, and wave/background are inherited from the native layout.
    add_accent(slide, 1.784)
    text_box(slide, 1.94, 0.69, 8.27, 1.58, spec.get("ending_title", "THANKS"), 78, PRIMARY_DEEP, False)
    lines = spec.get(
        "ending_lines",
        [
            "金风科技股份有限公司",
            "GOLDWIND SCIENCE & TECHNOLOGY CO., LTD",
            "北京金风科创风电设备有限公司",
            "BEIJING GOLDWIND SCIENCE & CREATION WINDPOWER EQUIPMENT CO., LTD",
        ],
    )
    y_positions = [2.70, 3.21, 4.01, 4.52]
    sizes = [19, 17, 19, 17]
    for text, y, size in zip(lines[:4], y_positions, sizes):
        text_box(slide, 1.95, y, 10.7, 0.42, text, size, BODY, size == 19)


def load_spec(path: Path | None) -> dict:
    if path:
        return json.loads(path.read_text(encoding="utf-8"))
    return {
        "title": "金风科技2025股东回报",
        "subtitle": "年度价值创造观察",
        "author": "汇报人 - Codex",
        "date": "2026-04-24",
        "toc": ["回报框架", "业绩兑现", "增长基础", "治理与ESG"],
        "pages": [
            {
                "title": "股东回报框架",
                "lead": "股东回报不是单点分红，而是经营质量、现金创造、分配纪律和透明沟通的连续闭环。",
                "bullets": [
                    {"heading": "经营质量", "body": "收入、利润、ROE决定长期回报的基本盘。"},
                    {"heading": "现金创造", "body": "经营现金流与货币资金决定分配安全边际。"},
                    {"heading": "增长基础", "body": "订单、全球化和服务容量支撑后续兑现。"},
                    {"heading": "治理沟通", "body": "信息披露和ESG实践影响长期信任。"},
                ],
                "source": "Goldwind 2025 Annual Results Report",
            },
            {
                "title": "2025年度业绩兑现",
                "lead": "2025年收入、归母净利润和ROE同步提升，股东回报的经营底座明显增强。",
                "bullets": [
                    {"heading": "营业收入 730.23亿", "body": "年度规模继续提升。"},
                    {"heading": "归母净利润 27.74亿", "body": "盈利恢复支撑回报能力。"},
                    {"heading": "综合毛利率 14.18%", "body": "业务质量持续修复。"},
                    {"heading": "加权ROE 7.08%", "body": "较2024年继续改善。"},
                ],
                "source": "Goldwind 2025 Annual Results Report",
            },
            {
                "title": "订单与全球化支撑",
                "lead": "订单、销售容量和国际化布局决定利润修复能否转化为跨周期回报能力。",
                "bullets": [
                    {"heading": "对外销售容量 26,626MW", "body": "同比增长65.9%。"},
                    {"heading": "在手订单 53.7GW", "body": "截至2025年底。"},
                    {"heading": "全球累计装机 165GW", "body": "业务基础持续扩张。"},
                    {"heading": "海外订单 9,270.17MW", "body": "国际业务提供增量线索。"},
                ],
                "source": "Goldwind 2025 Annual Results Report",
            },
            {
                "title": "现金质量与回报约束",
                "lead": "股东回报不能只看利润，还要看经营现金流、货币资金和已执行分红节奏。",
                "bullets": [
                    {"heading": "经营现金流 35.43亿", "body": "现金流恢复增强分配空间。"},
                    {"heading": "货币资金 103.23亿", "body": "年末资金安全垫。"},
                    {"heading": "已执行分红", "body": "2024年度末期股息每股0.14元。"},
                    {"heading": "后续观察", "body": "结合资本开支和订单交付节奏评估。"},
                ],
                "source": "Goldwind 2025 Annual Results Report",
            },
            {
                "title": "治理与ESG长期信任",
                "lead": "治理质量和ESG兑现会影响资本市场对长期现金流的稳定性判断。",
                "bullets": [
                    {"heading": "信息披露考核A", "body": "治理透明度保持较高水平。"},
                    {"heading": "绿色电力占比99%", "body": "全球生产及运营活动绿电使用。"},
                    {"heading": "供应链协同", "body": "主要零部件供应商绿电使用比例达100%。"},
                    {"heading": "投资者关系", "body": "持续强化业绩沟通与价值传递。"},
                ],
                "source": "Goldwind 2025 Annual Results Report",
            },
        ],
    }


def safe_output_name(spec: dict) -> str:
    title = str(spec.get("title") or "金风科技汇报").strip()
    cleaned = re.sub(r'[\\/:*?"<>|\r\n\t]+', "", title)
    cleaned = re.sub(r"\s+", "", cleaned)
    cleaned = cleaned[:36] or "金风科技汇报"
    return f"{cleaned}.pptx"


def build(spec: dict, output: Path, base: Path, toc_image: Path) -> Path:
    prs = Presentation(str(base))
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)
    cover_layout = get_layout(prs, "60_")
    content_layout = get_layout(prs, "8_")
    copyright_element = find_copyright_element(prs)
    clear_slides(prs)

    cover = prs.slides.add_slide(cover_layout)
    clone_shape_to_slide(cover, copyright_element)
    add_cover(cover, spec)

    toc = prs.slides.add_slide(content_layout)
    add_toc(toc, spec.get("toc") or [page.get("title", "") for page in spec.get("pages", [])[:4]], toc_image)

    pages = spec.get("pages") or []
    for idx, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(content_layout)
        add_content(slide, page, f"{idx:02d}")

    ending = prs.slides.add_slide(cover_layout)
    clone_shape_to_slide(ending, copyright_element)
    add_ending(ending, spec)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a native editable Goldwind-standard PPTX.")
    parser.add_argument("--spec", type=Path, help="JSON deck spec. If omitted, builds a bundled demo deck.")
    parser.add_argument("-o", "--output", type=Path, help="Output PPTX path")
    parser.add_argument("--base", type=Path, default=default_base())
    parser.add_argument("--toc-image", type=Path, default=default_toc_image())
    args = parser.parse_args()

    spec = load_spec(args.spec)
    output = args.output or Path.cwd() / safe_output_name(spec)
    build(spec, output, args.base, args.toc_image)
    print(f"Goldwind native deck built: {output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
