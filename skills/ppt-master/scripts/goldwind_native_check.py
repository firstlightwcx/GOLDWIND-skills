#!/usr/bin/env python3
"""Validate native Goldwind-standard PPTX decks.

This checker guards the fragile parts that previously failed in WPS/PowerPoint:
the TOC left image, native rotated side copyright, visible media relationships,
and correct Goldwind master/layout binding. It is intentionally structural so it
can run without opening Office.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from zipfile import ZipFile

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ModuleNotFoundError as exc:  # pragma: no cover - environment guard
    raise SystemExit("python-pptx is required. Install with: python3 -m pip install python-pptx") from exc


EMU_PER_INCH = 914400
EXPECTED_W = 13.333333
EXPECTED_H = 7.5
TOL = 0.08
COPYRIGHT = "GOLDWIND SCIENCE"


def inch(value: int) -> float:
    return value / EMU_PER_INCH


def near(value: float, expected: float, tol: float = TOL) -> bool:
    return abs(value - expected) <= tol


def text_of(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text.replace("\n", " ").strip()


def fail(errors: list[str], message: str) -> None:
    errors.append(message)


def check_zip_integrity(path: Path, errors: list[str]) -> None:
    with ZipFile(path) as zf:
        names = zf.namelist()
        slide_xml = [n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        media = [n for n in names if n.startswith("ppt/media/")]
        if not slide_xml:
            fail(errors, "PPTX has no slide XML files.")
        if not media:
            fail(errors, "PPTX has no embedded media; logo/background assets may be missing.")
        for item in zf.infolist():
            if item.file_size == 0 and item.filename.startswith(("ppt/slides/", "ppt/media/")):
                fail(errors, f"Empty PPTX part: {item.filename}")


def check_slide_size(prs: Presentation, errors: list[str]) -> None:
    w = inch(prs.slide_width)
    h = inch(prs.slide_height)
    if not (near(w, EXPECTED_W, 0.02) and near(h, EXPECTED_H, 0.02)):
        fail(errors, f"Unexpected slide size: {w:.3f} x {h:.3f} in; expected 13.333 x 7.500.")


def check_layouts(prs: Presentation, errors: list[str]) -> None:
    names = [layout.name for layout in prs.slide_layouts]
    if not any("60_" in name for name in names):
        fail(errors, "Missing Goldwind cover/ending layout containing '60_'.")
    if not any("8_" in name for name in names):
        fail(errors, "Missing Goldwind content/TOC layout containing '8_'.")
    for index, slide in enumerate(prs.slides, start=1):
        name = slide.slide_layout.name
        if index in (1, len(prs.slides)):
            if "60_" not in name:
                fail(errors, f"Slide {index} should use '60_' cover/ending layout, got {name!r}.")
        elif "8_" not in name:
            fail(errors, f"Slide {index} should use '8_' content/TOC layout, got {name!r}.")


def check_toc(prs: Presentation, errors: list[str]) -> None:
    if len(prs.slides) < 2:
        fail(errors, "Deck has fewer than 2 slides; TOC slide is missing.")
        return
    slide = prs.slides[1]
    texts = " ".join(text_of(shape) for shape in slide.shapes)
    if "目录" not in texts:
        fail(errors, "Slide 2 missing TOC title '目录'.")
    if "1." not in texts or "2." not in texts:
        fail(errors, "Slide 2 TOC items are missing numbered agenda entries.")
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and getattr(shape, "has_text_frame", False) and inch(shape.top) <= 1.0:
            fail(errors, "Slide 2 contains an unused top title placeholder; WPS may show placeholder prompt text.")
    left_pictures = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        x, y, w, h = inch(shape.left), inch(shape.top), inch(shape.width), inch(shape.height)
        if x < 0.25 and w > 6.0 and h > 7.0:
            left_pictures.append((x, y, w, h))
    if not left_pictures:
        fail(errors, "Slide 2 missing the left full-height TOC wind-turbine image.")
    else:
        x, y, w, h = left_pictures[0]
        if not (near(x, 0, 0.05) and near(y, -0.006, 0.08) and near(w, 6.92, 0.12) and near(h, 7.506, 0.12)):
            fail(errors, f"Slide 2 TOC left image anchor drifted: x={x:.3f}, y={y:.3f}, w={w:.3f}, h={h:.3f}.")


def copyright_shapes_for_slide(slide):
    for shape in slide.shapes:
        yield "slide", shape
    for shape in slide.slide_layout.shapes:
        yield "layout", shape


def check_horizontal_copyright(prs: Presentation, errors: list[str]) -> None:
    for slide_idx, slide in enumerate(prs.slides, start=1):
        found = False
        for scope, shape in copyright_shapes_for_slide(slide):
            txt = text_of(shape).upper()
            if COPYRIGHT not in txt or "©" not in text_of(shape):
                continue
            found = True
            rotation = float(getattr(shape, "rotation", 0) or 0)
            x = inch(shape.left)
            y = inch(shape.top)
            if abs(rotation - 270) > 1:
                fail(errors, f"Slide {slide_idx} has non-rotated {scope} copyright text at x={x:.3f}, y={y:.3f}.")
            if y < 0.35 and rotation == 0:
                fail(errors, f"Slide {slide_idx} {scope} copyright appears at top edge; side rail transform was lost.")
        if not found:
            fail(errors, f"Slide {slide_idx} missing Goldwind side copyright in slide or layout.")


def check_content_titles(prs: Presentation, errors: list[str]) -> None:
    for slide_zero_idx in range(2, len(prs.slides) - 1):
        index = slide_zero_idx + 1
        slide = prs.slides[slide_zero_idx]
        texts = [text_of(shape) for shape in slide.shapes if text_of(shape)]
        titles = [txt for txt in texts if "「" in txt and "」" in txt]
        if not titles:
            fail(errors, f"Slide {index} missing Goldwind content title format like 01「标题」.")
        if len(titles) > 1:
            fail(errors, f"Slide {index} has duplicate Goldwind content titles: {titles!r}.")
        top_placeholders = [
            text_of(shape)
            for shape in slide.shapes
            if getattr(shape, "is_placeholder", False) and getattr(shape, "has_text_frame", False) and inch(shape.top) <= 1.0
        ]
        if not any("「" in txt and "」" in txt for txt in top_placeholders):
            fail(errors, f"Slide {index} top title placeholder is not filled; WPS may show placeholder prompt text.")


def check(path: Path) -> list[str]:
    errors: list[str] = []
    if not path.exists():
        return [f"File not found: {path}"]
    check_zip_integrity(path, errors)
    prs = Presentation(str(path))
    if len(prs.slides) < 3:
        fail(errors, f"Deck should contain at least cover, TOC, and ending slides; got {len(prs.slides)}.")
        return errors
    check_slide_size(prs, errors)
    check_layouts(prs, errors)
    check_toc(prs, errors)
    check_horizontal_copyright(prs, errors)
    check_content_titles(prs, errors)
    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate a native editable Goldwind PPTX.")
    parser.add_argument("pptx", type=Path)
    args = parser.parse_args()
    errors = check(args.pptx)
    if errors:
        print("Goldwind native PPTX check FAILED", file=sys.stderr)
        for error in errors:
            print(f"- {error}", file=sys.stderr)
        return 1
    print("Goldwind native PPTX check PASSED")
    print(f"- file: {args.pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
