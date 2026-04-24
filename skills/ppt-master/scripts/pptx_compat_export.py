#!/usr/bin/env python3
"""Export SVG slides to a maximum-compatibility raster PPTX.

This is a fallback for WPS / older PowerPoint builds that may show blank slides
for pure SVG or partially supported DrawingML content. It renders each final SVG
with headless Chrome into a full-slide PNG, then places that PNG as the only
slide object in a PPTX. The output is not editable, but it is intentionally
robust for viewing and sharing.
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


DEFAULT_CHROME_PATHS = [
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
]


def find_chrome(explicit: str | None = None) -> str:
    candidates = [explicit] if explicit else []
    candidates.extend(DEFAULT_CHROME_PATHS)
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    for name in ("google-chrome", "chromium", "chromium-browser", "msedge"):
        resolved = shutil.which(name)
        if resolved:
            return resolved
    raise SystemExit(
        "No Chrome-compatible browser found. Install Google Chrome or pass --chrome /path/to/browser."
    )


def sorted_svgs(svg_dir: Path) -> list[Path]:
    svgs = sorted(svg_dir.glob("*.svg"))
    if not svgs:
        raise SystemExit(f"No SVG files found in {svg_dir}")
    return svgs


def render_svg(chrome: str, svg_path: Path, png_path: Path, width: int, height: int) -> None:
    png_path.parent.mkdir(parents=True, exist_ok=True)
    cmd = [
        chrome,
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--force-device-scale-factor=1",
        f"--window-size={width},{height}",
        f"--screenshot={png_path}",
        svg_path.resolve().as_uri(),
    ]
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
    if result.returncode != 0:
        raise SystemExit(f"Chrome render failed for {svg_path.name}:\n{result.stdout}")
    if not png_path.exists() or png_path.stat().st_size == 0:
        raise SystemExit(f"Chrome did not create a valid PNG for {svg_path.name}")


def create_pptx(pngs: list[Path], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    # Remove default first slide if any layout initialization creates one.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    parser = argparse.ArgumentParser(description="Export final SVG slides as a raster compatibility PPTX.")
    parser.add_argument("project", type=Path, help="PPT Master project directory")
    parser.add_argument("-s", "--source", default="final", choices=["output", "final"], help="SVG source directory")
    parser.add_argument("--chrome", help="Path to Google Chrome / Edge / Chromium executable")
    parser.add_argument("--width", type=int, default=1280)
    parser.add_argument("--height", type=int, default=720)
    parser.add_argument("-o", "--output", type=Path)
    args = parser.parse_args()

    project = args.project.resolve()
    svg_dir = project / ("svg_final" if args.source == "final" else "svg_output")
    svgs = sorted_svgs(svg_dir)
    chrome = find_chrome(args.chrome)

    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    render_dir = project / "images" / f"compat_render_{stamp}"
    pngs: list[Path] = []
    for idx, svg in enumerate(svgs, start=1):
        png = render_dir / f"{idx:02d}_{svg.stem}.png"
        render_svg(chrome, svg, png, args.width, args.height)
        pngs.append(png)

    output = args.output or project / "exports" / f"{project.name}_compat_{stamp}.pptx"
    create_pptx(pngs, output)
    print("Compatibility PPTX export complete")
    print(f"- source_svg_dir: {svg_dir}")
    print(f"- rendered_png_dir: {render_dir}")
    print(f"- slides: {len(pngs)}")
    print(f"- output: {output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
